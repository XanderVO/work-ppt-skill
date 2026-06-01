"""
build_deck.py
=============

Orchestrator script for the Bmatix PPT skill (v3 — two-template architecture).

There are TWO templates, selected based on whether the user has an existing
chat to work from:

  1. Empty_Bmatix_template.pptx — for FAST-PATH builds.
     19 slides: intro + TOC + 3 content placeholder slides + closing + 13
     hidden special-layout slides (timeline, world_map, etc.). Used when the
     user has no source chat. Output = the entire template, with intro and
     closing personalized; content slides stay as placeholders for the user
     to fill in PowerPoint themselves.

  2. Fill-in_Bmatix_template.pptx — for PAD-B builds (existing chat).
     3 slides: intro + 1 content canvas + closing. The content canvas is
     duplicated once per content slide in the plan, then cards/icons/pills
     are drawn programmatically on each duplicate using slide_helpers.

This script no longer supports "special layouts" (timeline, world_map, etc.)
as pad-B content slides. The only kind of content slide in a pad-B build is
"card". If users want a Gantt or timeline visual, they should pick fast-path
(which exposes the hidden special slides for them to use manually).

Usage:
  python build_deck.py \
    --template-empty path/to/Empty_Bmatix_template.pptx \
    --template-fillin path/to/Fill-in_Bmatix_template.pptx \
    --output path/to/output.pptx \
    --plan path/to/plan.json

Plan.json schema: see SKILL.md "Build the Deck" section.
"""

import argparse
import json
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

# Import slide helpers from the same scripts/ directory
sys.path.insert(0, str(Path(__file__).resolve().parent))
import slide_helpers as sh
from validate_plan import validate_plan, PlanValidationError


# =============================================================================
# TEMPLATE SLIDE INDICES (1-based)
# =============================================================================

# Empty template (fast-path)
EMPTY_INTRO_INDEX = 1
EMPTY_CLOSING_INDEX = 6  # "Questions?" end slide

# Fill-in template (pad B)
FILLIN_INTRO_INDEX = 1
FILLIN_CANVAS_INDEX = 2   # Content canvas to duplicate per content slide
FILLIN_CLOSING_INDEX = 3


# =============================================================================
# SLIDE MANIPULATION HELPERS
# =============================================================================

def _xml_slide_id_lst(prs):
    return prs.slides._sldIdLst


def _slide_id_entries(prs):
    return list(_xml_slide_id_lst(prs))


def remove_slide(prs, slide_index_zero_based):
    """Remove a slide by 0-based index. Manipulates XML directly since
    python-pptx doesn't expose a clean API."""
    sld_lst = _xml_slide_id_lst(prs)
    entries = _slide_id_entries(prs)
    entry = entries[slide_index_zero_based]
    rId = entry.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sld_lst.remove(entry)


def _remap_rid_references_in_xml(element, rid_remap: dict):
    """Replace any attribute value that matches an old rId with its new rId.

    Needed because deepcopy of a slide's shape XML keeps old rId strings
    literally, but the destination slide has its own relationship IDs.
    Without remapping, embedded images/icons break (r:embed="rId2" points
    at a relationship that doesn't exist on the destination).
    """
    if not rid_remap:
        return
    for el in element.iter():
        for attr_name, attr_val in list(el.attrib.items()):
            if attr_val in rid_remap:
                el.attrib[attr_name] = rid_remap[attr_val]


def duplicate_slide(prs, source_index_zero_based):
    """Duplicate a slide and append it to the end of the deck.

    Preserves:
      - Slide layout (so background/master inheritance is correct)
      - All shape XML (copied via deepcopy)
      - All relationships (images, hyperlinks, etc.) with remapped rIds
      - Explicit slide background if present

    Returns the new slide object.
    """
    slides = prs.slides
    source_slide = slides[source_index_zero_based]
    source_layout = source_slide.slide_layout

    new_slide = slides.add_slide(source_layout)

    # Remove default placeholders the layout added (we'll copy from source)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy relationships and build rId remap
    rid_remap = {}
    for rel_id, rel in source_slide.part.rels.items():
        reltype = rel.reltype
        if reltype.endswith("/slideLayout"):
            continue  # add_slide() already wired this up
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, reltype,
                                                is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, reltype)
        if new_rid != rel_id:
            rid_remap[rel_id] = new_rid

    # Copy each shape XML, remapping rIds in the copy
    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        _remap_rid_references_in_xml(new_el, rid_remap)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Copy explicit slide background if source has one
    src_bg = source_slide._element.find(qn("p:cSld")).find(qn("p:bg"))
    if src_bg is not None:
        dst_cSld = new_slide._element.find(qn("p:cSld"))
        existing_bg = dst_cSld.find(qn("p:bg"))
        if existing_bg is not None:
            dst_cSld.remove(existing_bg)
        new_bg = deepcopy(src_bg)
        _remap_rid_references_in_xml(new_bg, rid_remap)
        dst_cSld.insert(0, new_bg)

    return new_slide


def move_slide_to_position(prs, slide, target_index_zero_based):
    """Move a slide to a specific position in the deck order (0-based)."""
    sld_lst = _xml_slide_id_lst(prs)
    entries = _slide_id_entries(prs)
    target_rId = None
    for rel_id, rel in prs.part.rels.items():
        if rel.target_part is slide.part:
            target_rId = rel_id
            break
    if target_rId is None:
        raise RuntimeError("Could not find slide in part relationships")

    moved_entry = None
    for entry in entries:
        if entry.get(qn("r:id")) == target_rId:
            moved_entry = entry
            break
    if moved_entry is None:
        raise RuntimeError("Could not find slide in sldIdLst")

    sld_lst.remove(moved_entry)
    remaining = list(sld_lst)
    if target_index_zero_based >= len(remaining):
        sld_lst.append(moved_entry)
    else:
        sld_lst.insert(target_index_zero_based, moved_entry)


# =============================================================================
# MAIN BUILD ROUTINE
# =============================================================================

def build_deck(template_empty_path: Path, template_fillin_path: Path,
               output_path: Path, plan: dict, skip_validation: bool = False):
    """Build the deck per plan and save to output_path.

    Branches:
      - If `plan['content_slides']` is empty (or missing): FAST-PATH. Open
        the Empty template, fill intro and closing only, save as-is.
      - Otherwise: PAD B. Open the Fill-in template, fill intro, duplicate
        the canvas slide once per content slide, draw cards on each
        duplicate, fill closing, remove the original canvas, save.

    Validates plan.json against hard rules BEFORE building. On violation,
    raises PlanValidationError with the full list of issues. Set
    `skip_validation=True` only for emergency debug builds (NOT recommended).
    """
    intro = plan.get("intro", {})
    closing = plan.get("closing", {})
    content_specs = plan.get("content_slides", [])
    is_fastpath = (len(content_specs) == 0)

    # Hard-rule validation (raises PlanValidationError if anything fails).
    # Fast-path plans skip validation by design (the Empty template is
    # delivered as-is, with no per-slide content to validate).
    if not skip_validation and not is_fastpath:
        validate_plan(plan)

    if is_fastpath:
        _build_fastpath(template_empty_path, output_path, intro, closing)
    else:
        _build_pad_b(template_fillin_path, output_path, intro, closing,
                     content_specs)


def _build_fastpath(template_path: Path, output_path: Path,
                    intro: dict, closing: dict):
    """Fast-path build: deliver the Empty template integrally with only the
    intro and closing personalized. Content placeholder slides (2-5) and
    hidden special-layout slides (7-19) stay as-is for the user to fill in
    PowerPoint themselves.
    """
    prs = Presentation(str(template_path))

    # Fill intro (slide 1)
    intro_slide = prs.slides[EMPTY_INTRO_INDEX - 1]
    sh.fill_intro_slide(
        intro_slide,
        title=intro.get("title", ""),
        subtitle=intro.get("subtitle", ""),
        author_date=intro.get("author_date", ""),
    )

    # Fill closing (slide 6)
    closing_slide = prs.slides[EMPTY_CLOSING_INDEX - 1]
    sh.fill_closing_slide(
        closing_slide,
        author_date=closing.get("author_date", intro.get("author_date", "")),
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _build_pad_b(template_path: Path, output_path: Path,
                 intro: dict, closing: dict, content_specs: list):
    """Pad-B build: open the Fill-in template, fill intro/closing, duplicate
    the content canvas once per content_spec, draw cards on each duplicate,
    remove the original canvas.
    """
    prs = Presentation(str(template_path))

    # Fill intro (slide 1)
    intro_slide = prs.slides[FILLIN_INTRO_INDEX - 1]
    sh.fill_intro_slide(
        intro_slide,
        title=intro.get("title", ""),
        subtitle=intro.get("subtitle", ""),
        author_date=intro.get("author_date", ""),
    )

    # Get a reference to the original canvas slide BEFORE we start
    # duplicating; we'll remove it at the end.
    canvas_slide = prs.slides[FILLIN_CANVAS_INDEX - 1]
    closing_slide = prs.slides[FILLIN_CLOSING_INDEX - 1]

    # Fill closing first (we know its index before things move around)
    sh.fill_closing_slide(
        closing_slide,
        author_date=closing.get("author_date", intro.get("author_date", "")),
    )

    # For each content slide spec, duplicate the canvas, fill it, and move
    # it to the correct position.
    # The duplicate ends up at the end of the deck; we'll move it into place
    # just before the closing slide.
    new_content_slides = []
    canvas_idx_zero_based = FILLIN_CANVAS_INDEX - 1
    for spec in content_specs:
        kind = spec.get("kind", "card")
        if kind != "card":
            raise ValueError(
                f"Pad-B builds only support kind='card'. Got '{kind}'. "
                f"Special layouts are fast-path-only in this version."
            )
        new_slide = duplicate_slide(prs, canvas_idx_zero_based)
        sh.draw_card_slide(new_slide, spec)
        new_content_slides.append(new_slide)

    # Reorder: [intro, new_content_slides..., closing, canvas_template].
    # We move each new content slide into position 1, 2, 3... (just after intro).
    # The closing slide stays where it is until we move the canvas.
    # Approach: move new slides one by one to position right before closing.

    # We move slides into their correct positions left-to-right.
    # Position 0 = intro. Position 1 = first content. Position N = last content.
    # Position N+1 = closing. Position N+2 = leftover canvas (to be deleted).
    target_pos = 1  # right after intro
    for new_slide in new_content_slides:
        move_slide_to_position(prs, new_slide, target_pos)
        target_pos += 1

    # Now move closing to right after the last content slide
    move_slide_to_position(prs, closing_slide, target_pos)

    # The original canvas slide should now be the last one. Remove it.
    # We re-fetch the slide list to find its current index.
    canvas_position = None
    for i, slide in enumerate(prs.slides):
        if slide is canvas_slide:
            canvas_position = i
            break
    if canvas_position is None:
        raise RuntimeError("Lost track of the original canvas slide; cannot remove it.")
    remove_slide(prs, canvas_position)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# =============================================================================
# CLI ENTRY POINT
# =============================================================================

def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template-empty", required=True, type=Path,
                        help="Path to Empty_Bmatix_template.pptx (fast-path)")
    parser.add_argument("--template-fillin", required=True, type=Path,
                        help="Path to Fill-in_Bmatix_template.pptx (pad B)")
    parser.add_argument("--output", required=True, type=Path,
                        help="Output .pptx path")
    parser.add_argument("--plan", required=True, type=Path,
                        help="Path to plan.json")
    parser.add_argument("--skip-validation", action="store_true",
                        help="Skip hard-rule validation (emergency debug only)")
    args = parser.parse_args()

    plan = json.loads(args.plan.read_text(encoding="utf-8"))
    try:
        build_deck(args.template_empty, args.template_fillin, args.output,
                   plan, skip_validation=args.skip_validation)
    except PlanValidationError as e:
        # Print one issue per line to stderr in a parseable format so Claude
        # can read each issue independently and auto-fix.
        print(f"VALIDATION FAILED: {len(e.issues)} issue(s)", file=sys.stderr)
        for issue in e.issues:
            print(str(issue), file=sys.stderr)
        sys.exit(2)
    print(f"Deck written to {args.output}")


if __name__ == "__main__":
    main()
