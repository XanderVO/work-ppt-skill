"""
slide_helpers.py
================

All primitives + card-drawing functions for the Bmatix PPT skill.

Architecture
------------
- Constants block (colors, sizes, fonts) at the top - single source of truth.
- Primitive helpers (cards, pills, circles, arrows, image placeholders).
- Shared header/footer drawer used by all card layouts.
- 8 card-layout drawers, one per layout key in the Card Catalog.
- Intro/closing placeholder fillers for the Fill-in template.

The build_deck.py orchestrator calls these functions. This file does NOT
open/save the .pptx; it only mutates a Presentation object passed in.

Color tokens align with the Bmatix master theming (see SKILL.md
"Bmatix Brand Colors" for the full palette).
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy


# =============================================================================
# CONSTANTS  ----  edit here, propagates everywhere
# =============================================================================

# --- Brand colors (hex, no #) ----------------------------------------------
# These match the Claude_Cards_template.pptx exactly (see SKILL.md for color
# lineage). The template is the source of truth.
NAVY        = RGBColor(0x15, 0x1F, 0x6D)   # primary navy (template-matched)
ORANGE      = RGBColor(0xFF, 0x69, 0x00)   # primary orange (template-matched)
LIGHT_BLUE  = RGBColor(0x54, 0xB0, 0xE7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# Card fills + accent stripes (template-matched)
CARD_FILL         = RGBColor(0xEE, 0xF0, 0xFA)  # lavender-blue (default)
CARD_FILL_EXAMPLE = RGBColor(0xEF, 0xE8, 0xDD)  # warm beige (for example/demo cards)
CARD_FILL_LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)  # near-white (image placeholder bg)
STRIPE_ORANGE     = ORANGE
STRIPE_NAVY       = NAVY
STRIPE_GREEN      = RGBColor(0x3F, 0x8F, 0x3A)
STRIPE_RED        = RGBColor(0xC2, 0x4A, 0x4A)
STRIPE_EXAMPLE    = RGBColor(0x8C, 0x7B, 0x62)  # tan/brown for example card stripe

# Map for plan.json "stripe": "orange|navy|green|red|example"
STRIPE_COLORS = {
    "orange":  STRIPE_ORANGE,
    "navy":    STRIPE_NAVY,
    "green":   STRIPE_GREEN,
    "red":     STRIPE_RED,
    "example": STRIPE_EXAMPLE,
}

# --- Slide geometry (inches) -----------------------------------------------
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.52
MARGIN_R = 0.52
MARGIN_T = 0.36
MARGIN_B = 0.30

# Header zone: title + subtitle + use case bar
TITLE_TOP      = 0.36
TITLE_HEIGHT   = 0.77
SUBTITLE_TOP   = 1.26
SUBTITLE_HEIGHT = 0.38
USECASE_TOP    = 1.70
USECASE_HEIGHT = 0.30

# Content zone (below header, above footer)
CONTENT_TOP    = 2.15
CONTENT_BOTTOM = 7.00     # leave 0.5" for footer
CONTENT_LEFT   = MARGIN_L
CONTENT_RIGHT  = SLIDE_W - MARGIN_R

# Footer
FOOTER_Y       = 7.05

# --- Card defaults ---------------------------------------------------------
CARD_CORNER_RADIUS = 0.04   # in EMU fraction (python-pptx uses adjustment 0-1)
CARD_PADDING       = 0.20   # interior padding inches
CARD_STRIPE_W      = 0.07   # side stripe width inches (uniform across all cards)
CARD_GAP           = 0.25   # gap between cards inches

# --- Typography (points) ---------------------------------------------------
FONT_FAMILY        = "Calibri"
FONT_TITLE_SIZE    = 32
FONT_SUBTITLE_SIZE = 18
FONT_USECASE_SIZE  = 12
FONT_CARD_LABEL    = 14   # was 11 — labels need to be readable at presentation distance
FONT_CARD_TITLE    = 20   # was 14 — card titles need to be readable at presentation distance
FONT_CARD_BODY     = 16   # was 11 — card bodies need to be readable at presentation distance
FONT_PILL          = 9
FONT_FOOTER        = 9
FONT_NUMBER        = 16
FONT_ICON_LABEL    = 10

# Underline thickness (points) for orange subtitle line + navy usecase line
UNDERLINE_PT = 1.5
UNDERLINE_LEN_IN = 1.4      # short accent line under subtitle/usecase
USECASE_UNDERLINE_LEN_IN = 0.73   # matches width of "USE CASE" text at 12pt bold

# --- Footer text -----------------------------------------------------------
FOOTER_COPYRIGHT   = "© Bmatix 2026"
FOOTER_ACT_INFORMED = "Act informed"


# =============================================================================
# LOW-LEVEL PRIMITIVES
# =============================================================================

def _set_solid_fill(shape, rgb: RGBColor):
    """Apply a solid fill of the given color to a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_no_fill(shape):
    shape.fill.background()


def _set_no_line(shape):
    shape.line.fill.background()


def _set_line(shape, rgb: RGBColor, width_pt: float = 0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def _set_text(shape, text: str, *, font_size: int = 12,
              bold: bool = False, color: RGBColor = NAVY,
              align: PP_ALIGN = PP_ALIGN.LEFT, font_name: str = FONT_FAMILY,
              anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
              italic: bool = False):
    """
    Set text on a shape's text_frame with formatting.

    Single-paragraph helper. For multi-line text, pass '\\n' in `text` and
    each line becomes a paragraph with the same formatting.
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # Reduce internal margins so text uses the full shape area
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    lines = text.split("\n")
    # Clear existing paragraphs by overwriting first and adding rest
    tf.paragraphs[0].text = ""
    for run in tf.paragraphs[0].runs:
        run.text = ""

    p0 = tf.paragraphs[0]
    p0.alignment = align
    r0 = p0.add_run()
    r0.text = lines[0]
    r0.font.name = font_name
    r0.font.size = Pt(font_size)
    r0.font.bold = bold
    r0.font.italic = italic
    r0.font.color.rgb = color

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color


def draw_rectangle(slide, x_in, y_in, w_in, h_in, *,
                   fill: RGBColor = None, line: RGBColor = None,
                   line_width_pt: float = 0.75,
                   rounded: bool = False):
    """Generic rectangle. Returns the shape."""
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_kind, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    if rounded:
        # Set corner adjustment to a small value (default is 0.35 - too round)
        shape.adjustments[0] = CARD_CORNER_RADIUS

    if fill is None:
        _set_no_fill(shape)
    else:
        _set_solid_fill(shape, fill)

    if line is None:
        _set_no_line(shape)
    else:
        _set_line(shape, line, line_width_pt)

    return shape


def draw_text_box(slide, x_in, y_in, w_in, h_in, text: str, **kwargs):
    """Plain textbox without fill. Returns the shape."""
    tb = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    _set_text(tb, text, **kwargs)
    return tb


def draw_line(slide, x1_in, y1_in, x2_in, y2_in,
              color: RGBColor = NAVY, width_pt: float = 1.0):
    """Draw a straight line connector. Returns the shape."""
    line = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1_in), Inches(y1_in), Inches(x2_in), Inches(y2_in)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


# =============================================================================
# CARD PRIMITIVES (cards, pills, circles, arrows, image placeholders)
# =============================================================================

def draw_card(slide, x_in, y_in, w_in, h_in, *,
              stripe_color: RGBColor = None,
              fill: RGBColor = CARD_FILL,
              stripe_width_in: float = None):
    """
    Draw a card: rounded rectangle with optional left side stripe.

    Returns dict with keys:
      'card'    - the card rectangle shape
      'stripe'  - the stripe shape, or None
      'inner'   - tuple (x, y, w, h) in inches for content placement,
                  accounting for stripe + padding

    All cards share the SAME stripe width (CARD_STRIPE_W) for visual
    consistency across slides. Result/conclusion cards are signaled by the
    stripe COLOR (orange) and the title color (orange), not by stripe width.
    """
    if stripe_width_in is None:
        stripe_width_in = CARD_STRIPE_W

    card = draw_rectangle(slide, x_in, y_in, w_in, h_in,
                          fill=fill, line=None, rounded=True)
    stripe = None
    inner_x = x_in + CARD_PADDING
    if stripe_color is not None:
        stripe = draw_rectangle(slide, x_in, y_in, stripe_width_in, h_in,
                                fill=stripe_color, line=None, rounded=False)
        inner_x = x_in + stripe_width_in + CARD_PADDING

    inner_y = y_in + CARD_PADDING
    inner_w = (x_in + w_in - CARD_PADDING) - inner_x
    inner_h = h_in - 2 * CARD_PADDING
    return {
        "card": card, "stripe": stripe,
        "inner": (inner_x, inner_y, inner_w, inner_h)
    }


# Default fixed height of the filled header band on a "filled header" card.
# Tuned so a 0.40" diameter icon-circle fits comfortably with vertical padding.
FILLED_HEADER_H = 0.55


def draw_card_with_filled_header(slide, x_in, y_in, w_in, h_in, *,
                                 header_color: RGBColor = NAVY,
                                 header_text: str = "",
                                 header_text_color: RGBColor = WHITE,
                                 icon: str = "",
                                 icon_font: str = None,
                                 icon_circle_color: RGBColor = None,
                                 body_fill: RGBColor = CARD_FILL,
                                 header_h: float = FILLED_HEADER_H,
                                 header_font_size: int = FONT_CARD_TITLE):
    """
    Card with a solid-color header band on top + lavender (or other) body
    below. Icon-circle (optional) sits inside the header band, on the left,
    with the title text to its right.

    Visual structure:
      ┌─────────────────────────────────────────┐
      │ [●]  HEADER TITLE                       │   <- header_color band
      ├─────────────────────────────────────────┤
      │ Body text in body_fill colored area     │   <- lavender by default
      │                                         │
      └─────────────────────────────────────────┘

    Returns dict with keys:
      'card'    - the body rectangle (the larger card shape underneath)
      'header'  - the header band rectangle
      'circle'  - the icon circle, or None
      'inner'   - (x, y, w, h) tuple in inches for body content placement
                  (below the header, inside body padding)

    Notes:
    - The card is drawn as a SINGLE rounded rectangle for the body + a
      flat rectangle for the header. The header sits flush on top so
      visually it looks like one card with a top band. Rounding the top
      corners cleanly across two shapes isn't worth the XML complexity
      here; the flat top-edge look is intentional and matches modern
      Bmatix-style decks.
    - When `icon` is empty, the title left-aligns inside the header with
      the standard left padding.
    """
    if icon_font is None:
        icon_font = FONT_FAMILY
    if icon_circle_color is None:
        # Default: white circle on navy header. White circle stands out
        # against the dark band and matches the white header text.
        icon_circle_color = WHITE

    # Body card (rounded). We let it cover the full height; the header
    # rectangle then sits on top and visually replaces the top portion.
    body = draw_rectangle(slide, x_in, y_in, w_in, h_in,
                          fill=body_fill, line=None, rounded=True)

    # Header band (flat, sits on top of body's top edge).
    header = draw_rectangle(slide, x_in, y_in, w_in, header_h,
                            fill=header_color, line=None, rounded=False)

    # Optional icon circle on the left side of the header.
    circle = None
    circle_d = header_h * 0.78  # 78% of header height (was 65% — bumped for emoji visibility)
    text_x = x_in + CARD_PADDING
    if icon:
        cx_in = x_in + CARD_PADDING + circle_d / 2
        cy_in = y_in + header_h / 2
        # Icon glyph color = the header's color (so it "punches through"
        # the white circle and visually connects to the band).
        circle = draw_circle(slide, cx_in, cy_in, circle_d,
                             fill=icon_circle_color, text=icon,
                             text_color=header_color,
                             font_size=int(circle_d * 32),
                             font_name=icon_font)
        text_x = cx_in + circle_d / 2 + 0.10

    # Header title text (vertically centered in band, to the right of icon).
    if header_text:
        text_w = (x_in + w_in) - text_x - CARD_PADDING
        tb = slide.shapes.add_textbox(
            Inches(text_x), Inches(y_in),
            Inches(text_w), Inches(header_h)
        )
        _set_text(tb, header_text,
                  font_size=header_font_size, bold=True,
                  color=header_text_color, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.MIDDLE)

    # Inner content area: below the header, inside body padding.
    inner_x = x_in + CARD_PADDING
    inner_y = y_in + header_h + 0.10
    inner_w = w_in - 2 * CARD_PADDING
    inner_h = h_in - header_h - 0.10 - CARD_PADDING

    return {
        "card": body, "header": header, "circle": circle,
        "inner": (inner_x, inner_y, inner_w, inner_h),
    }


def _is_example_card(card_spec):
    """
    Detect whether a card should be styled as an example/demo card (beige fill).

    Strict rule: ONLY when the word "example" or "voorbeeld" appears literally
    in the card's label, title, or header (case-insensitive). This avoids
    false positives - real content cards stay lavender.

    Examples that trigger beige fill:
      - {"label": "Example: vague prompt"} -> beige
      - {"title": "Voorbeeld"} -> beige
      - {"header": "Example output"} -> beige
    Cards that do NOT trigger beige (stay lavender):
      - {"label": "Vague prompt"}
      - {"title": "Claude's output"}
      - {"label": "The prompt"}
    """
    keywords = ("example", "voorbeeld")
    for key in ("label", "title", "header"):
        val = card_spec.get(key, "")
        if not isinstance(val, str):
            continue
        val_lower = val.lower()
        if any(kw in val_lower for kw in keywords):
            return True
    return False


def _resolve_card_fill(card_spec):
    """
    Pick card fill color based on spec.

    Default = lavender. Beige is used ONLY when the word "example" or
    "voorbeeld" appears literally in the card's label/title/header (see
    `_is_example_card`).
    """
    if _is_example_card(card_spec):
        return CARD_FILL_EXAMPLE
    return CARD_FILL


def _is_result_card(card_spec):
    """
    Detect whether a card should be styled as a result/conclusion card.

    A card is a result/conclusion/lesson-learned card if EITHER:
    - `is_result` is explicitly True
    - the label/title is one of the known result-like words (case-insensitive),
      including "result", "conclusion", "outcome", "summary", "takeaway",
      "lesson learned", "lessons learned", "key learnings", "insights",
      "key insights", "main takeaway", "key takeaway"

    Result/conclusion/lesson cards get:
    - orange side stripe
    - orange card title (instead of the default navy)
    Matches the "CONCLUSION" ribbon on slide 4 of the Claude Cards template.
    """
    if card_spec.get("is_result"):
        return True
    for key in ("label", "title", "header"):
        val = card_spec.get(key, "")
        if not isinstance(val, str):
            continue
        if val.strip().lower() in {"result", "results", "conclusion",
                                   "outcome", "summary", "key takeaway",
                                   "takeaway", "main takeaway",
                                   "lesson learned", "lessons learned",
                                   "key learnings", "learnings",
                                   "insights", "key insights"}:
            return True
    return False


def _resolve_stripe_for_card(card_spec):
    """
    Resolve the stripe color for a card.

    Returns a tuple (stripe_color or None, stripe_width_in) for backwards
    compatibility with callers. Stripe width is uniform across ALL cards
    (no thick/thin variation), so the second element is always CARD_STRIPE_W.

    Logic (priority order):
    1. If the card is a result/conclusion/lesson-learned card (per
       `_is_result_card`), force an ORANGE stripe — the visual signal that
       it's the takeaway / summary / conclusion. The card title also turns
       orange. Matches slide 4 (conclusion ribbon) in the Claude Cards
       template.
    2. If the card is an example card (per `_is_example_card`), force the
       tan/brown EXAMPLE stripe. Matches slide 7 (EXAMPLE · Claude Team set-up).
    3. Otherwise, honor the explicit `stripe` field if provided.
    """
    if _is_result_card(card_spec):
        return STRIPE_ORANGE, CARD_STRIPE_W

    if _is_example_card(card_spec):
        return STRIPE_EXAMPLE, CARD_STRIPE_W

    s = card_spec.get("stripe")
    if not s:
        # DEFAULT: all regular cards get a navy stripe so the thick side-bar
        # design element is visible deck-wide, not only on result/example
        # cards. Callers can suppress with `stripe: "none"`.
        return NAVY, CARD_STRIPE_W
    if s.lower() == "none":
        return None, CARD_STRIPE_W
    color = STRIPE_COLORS.get(s.lower())
    return color, CARD_STRIPE_W


def draw_pill_badge(slide, x_in, y_in, w_in, h_in, text: str,
                    *, fill: RGBColor = NAVY, text_color: RGBColor = WHITE):
    """Small pill badge (rounded rectangle with text). Returns shape."""
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    pill.adjustments[0] = 0.5  # fully rounded
    _set_solid_fill(pill, fill)
    _set_no_line(pill)
    _set_text(pill, text, font_size=FONT_PILL, bold=True,
              color=text_color, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    return pill


def draw_circle(slide, cx_in, cy_in, diameter_in, *,
                fill: RGBColor = NAVY, text: str = "",
                text_color: RGBColor = WHITE, font_size: int = FONT_NUMBER,
                bold: bool = True, font_name: str = FONT_FAMILY):
    """Draw a filled circle centered on (cx, cy). Returns shape.

    `text` can be either a number ("1", "2"...) or a Unicode/font-glyph icon
    (e.g. "✓", "⚡", "→", "★", or a Segoe MDL2/Wingdings glyph). For glyphs
    that require a specific font to render correctly, pass `font_name` (e.g.
    "Segoe MDL2 Assets" or "Wingdings").
    """
    x = cx_in - diameter_in / 2
    y = cy_in - diameter_in / 2
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y),
        Inches(diameter_in), Inches(diameter_in)
    )
    _set_solid_fill(circ, fill)
    _set_no_line(circ)
    if text:
        _set_text(circ, text, font_size=font_size, bold=bold,
                  color=text_color, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, font_name=font_name)
    return circ


# --- Icon-circle color palette --------------------------------------------
# When a step/subcard spec sets `circle_color`, map it through this table.
# Default stays NAVY (Bmatix primary). Other colors are reserved for
# "highlight" steps / "warning" steps / etc.
CIRCLE_COLORS = {
    "navy":   NAVY,
    "orange": ORANGE,
    "green":  STRIPE_GREEN,
    "red":    STRIPE_RED,
}


def resolve_circle_color(spec: dict) -> RGBColor:
    """Read `circle_color` from a step/subcard spec, fall back to NAVY."""
    key = (spec or {}).get("circle_color", "navy")
    return CIRCLE_COLORS.get(str(key).lower(), NAVY)


def draw_arrow_right(slide, x_in, y_in, w_in, h_in,
                     color: RGBColor = NAVY):
    """Right-pointing arrow shape (for between numbered steps)."""
    arr = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    _set_solid_fill(arr, color)
    _set_no_line(arr)
    return arr


def draw_image_placeholder(slide, x_in, y_in, w_in, h_in,
                           caption: str = "Image / mockup placeholder"):
    """
    Dashed-border rectangle that hints where the user should drop an image.
    Caption sits centered inside.
    """
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    _set_solid_fill(rect, CARD_FILL_LIGHT)
    rect.line.color.rgb = NAVY
    rect.line.width = Pt(1.0)
    # Dashed line via XML manipulation
    ln = rect.line._get_or_add_ln()
    prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    # Remove any existing prstDash
    for existing in ln.findall(qn("a:prstDash")):
        ln.remove(existing)
    ln.append(prstDash)

    _set_text(rect, f"🖼  {caption}",
              font_size=FONT_ICON_LABEL, bold=False, italic=True,
              color=NAVY, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    return rect


# =============================================================================
# =============================================================================
# SHARED HEADER + SUBTITLE FILLER (used by all card layouts)
# =============================================================================
#
# In v3 (two-template architecture) card slides are drawn on top of the
# Fill-in template's native content canvas (slide 2). That canvas already
# has TITLE and BODY/SUBTITLE placeholders styled by the Bmatix master.
# So all the card layouts just call draw_slide_chrome() to populate the
# native placeholders — they no longer draw their own title textboxes or
# the orange arc that used to live under the title in earlier versions.
#
# The legacy helpers (_estimate_title_width_in, _add_bmatix_arc,
# draw_slide_header, draw_slide_footer) were removed when the template
# architecture changed; the native master handles the title styling now.


def draw_slide_chrome(slide, *, title: str, subtitle: str = "",
                      use_case: str = ""):
    """Fill the slide's native TITLE and SUBTITLE placeholders.

    In the new template architecture (v3), card slides are NOT drawn on a
    blank white rectangle anymore. They are drawn on a duplicate of the
    Fill-in template's content canvas (slide 2). That canvas already has
    native title + subtitle placeholders styled per the Bmatix master
    (navy bold centered title, orange centered subtitle).

    So instead of drawing our own title textbox + orange arc + subtitle
    textbox like the old draw_slide_header() did, we just write into the
    placeholders that are already there. The native master/layout handles
    fonts, colors, and centering.

    `use_case` is accepted for backward compatibility with the card layout
    functions but is no longer rendered. The old "USE CASE · ..." line and
    its underline stripe relied on hand-drawn shapes; in the new
    architecture the simplified subtitle placeholder is the only secondary
    line under the title. If a caller passes a use_case, it gets appended
    to the subtitle (separated by " · ") so the information isn't lost.
    """
    # Compose subtitle + optional use_case
    sub_parts = []
    if subtitle:
        sub_parts.append(subtitle)
    if use_case:
        sub_parts.append(f"USE CASE · {use_case}")
    composed_subtitle = "  ·  ".join(sub_parts) if sub_parts else ""

    if title:
        # Card-canvas slides use the regular TITLE placeholder (not CENTER_TITLE)
        ok = _fill_placeholder_by_type(slide, "TITLE (1)", title)
        if not ok:
            # Fallback: try CENTER_TITLE if the layout happens to use that
            _fill_placeholder_by_type(slide, "CENTER_TITLE (3)", title)

    if composed_subtitle:
        # The Fill-in canvas slide has a BODY placeholder serving as subtitle
        # (positioned just below the title). Try BODY first, then SUBTITLE.
        ok = _fill_placeholder_by_type(slide, "BODY (2)", composed_subtitle)
        if not ok:
            _fill_placeholder_by_type(slide, "SUBTITLE (4)", composed_subtitle)


# =============================================================================
# CARD LAYOUTS (path a - 6 layouts)
# =============================================================================

def _card_title_color(card_spec):
    """
    Resolve title color from spec.

    Rules:
    - Explicit `title_color: "orange"` -> ORANGE
    - Result/conclusion card (detected via `_is_result_card`) -> ORANGE
    - Otherwise -> NAVY
    """
    if card_spec.get("title_color", "").lower() == "orange":
        return ORANGE
    if _is_result_card(card_spec):
        return ORANGE
    return NAVY


def _resolve_stripe(card_spec):
    """
    Backwards-compatible helper: returns stripe color only (no width).

    Most card-layout functions now call `_resolve_stripe_for_card` instead,
    which returns (color, width). This function remains for any callers
    that only need the color.
    """
    color, _ = _resolve_stripe_for_card(card_spec)
    return color


# ---- cards_2x2_compare -----------------------------------------------------

def draw_cards_2x2_compare(slide, slide_data):
    """
    2x2 grid of 4 cards. Each card has optional left stripe + label header + body.

    slide_data keys:
      title, subtitle, use_case  (passed to header)
      cards: list of 4 dicts with keys
        label   (e.g. "VAGUE PROMPT", bold uppercase)
        body    (multi-line text)
        stripe  ("red" | "green" | "orange" | "navy" | None)
        title_color ("orange" | "navy" | None) - affects label color
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    cards = slide_data.get("cards", [])
    # Pad to 4 if short, truncate if long
    cards = (cards + [{}] * 4)[:4]

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    card_w = (area_w - CARD_GAP) / 2
    card_h = (area_h - CARD_GAP) / 2

    positions = [
        (CONTENT_LEFT,             area_top),                          # top-left
        (CONTENT_LEFT + card_w + CARD_GAP, area_top),                  # top-right
        (CONTENT_LEFT,             area_top + card_h + CARD_GAP),      # bot-left
        (CONTENT_LEFT + card_w + CARD_GAP, area_top + card_h + CARD_GAP),  # bot-right
    ]

    for spec, (x, y) in zip(cards, positions):
        stripe_color, stripe_w = _resolve_stripe_for_card(spec)
        card_fill = _resolve_card_fill(spec)
        card = draw_card(slide, x, y, card_w, card_h,
                         stripe_color=stripe_color,
                         fill=card_fill,
                         stripe_width_in=stripe_w)
        ix, iy, iw, ih = card["inner"]

        # Label header
        label = spec.get("label", "")
        if label:
            label_color = _card_title_color(spec)
            draw_text_box(slide, ix, iy, iw, 0.35,
                          label.upper(),
                          font_size=FONT_CARD_LABEL, bold=True,
                          color=label_color, align=PP_ALIGN.LEFT)

        # Body
        body = spec.get("body", "")
        if body:
            body_y = iy + 0.40
            body_h = ih - 0.40
            draw_text_box(slide, ix, body_y, iw, body_h,
                          body,
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.LEFT)


# ---- cards_left_image_right ------------------------------------------------

def draw_cards_left_image_right(slide, slide_data):
    """
    3 stacked cards on left (icon + title + body) + image placeholder on right.

    slide_data keys:
      title, subtitle, use_case
      cards: list of 3 dicts with:
        icon (string emoji or short text)
        title
        body
        stripe (optional - typically orange for "Result" card)
      image_caption (optional)
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    cards = (slide_data.get("cards", []) + [{}] * 3)[:3]
    image_caption = slide_data.get("image_caption", "Image / mockup placeholder")

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    left_w = area_w * 0.42
    right_w = area_w * 0.55
    right_x = CONTENT_LEFT + left_w + (area_w - left_w - right_w)

    # 3 stacked cards on left
    card_h = (area_h - 2 * CARD_GAP) / 3
    for i, spec in enumerate(cards):
        cy = area_top + i * (card_h + CARD_GAP)
        stripe_color, stripe_w = _resolve_stripe_for_card(spec)
        card_fill = _resolve_card_fill(spec)
        card = draw_card(slide, CONTENT_LEFT, cy, left_w, card_h,
                         stripe_color=stripe_color,
                         fill=card_fill,
                         stripe_width_in=stripe_w)
        ix, iy, iw, ih = card["inner"]

        # Icon circle on left (small). Caller can vary `icon` (default "●"),
        # `circle_color` (navy/orange/green/red), and `icon_font` (default
        # Calibri; use "Segoe MDL2 Assets" or "Wingdings" for native PowerPoint
        # icons — those won't render in LibreOffice QA preview).
        icon = spec.get("icon", "●")
        icon_font = spec.get("icon_font", FONT_FAMILY)
        circle_fill = resolve_circle_color(spec)
        icon_d = 0.50
        draw_circle(slide, ix + icon_d / 2, iy + icon_d / 2 + 0.05,
                    icon_d, fill=circle_fill, text=icon,
                    text_color=WHITE, font_size=18, font_name=icon_font)

        # Title to the right of icon
        title_color = _card_title_color(spec)
        title_x = ix + icon_d + 0.15
        title_w = iw - icon_d - 0.15
        draw_text_box(slide, title_x, iy, title_w, 0.30,
                      spec.get("title", ""),
                      font_size=FONT_CARD_TITLE, bold=True,
                      color=title_color, align=PP_ALIGN.LEFT)

        # Body below title
        body = spec.get("body", "")
        if body:
            draw_text_box(slide, title_x, iy + 0.35, title_w, ih - 0.35,
                          body, font_size=FONT_CARD_BODY,
                          color=NAVY, align=PP_ALIGN.LEFT)

    # Image placeholder on right
    draw_image_placeholder(slide, right_x, area_top, right_w, area_h,
                           caption=image_caption)


# ---- cards_numbered_steps --------------------------------------------------

def draw_cards_numbered_steps(slide, slide_data):
    """
    3-4 numbered step cards horizontally + arrows between + 1 wide conclusion card.

    slide_data keys:
      title, subtitle, use_case
      header_style: "external" (default) places the icon-circle ABOVE each card.
                    "filled" gives each step card a navy/orange top band with
                    the icon-circle LEFT of the step title (white text on dark
                    band), and a lavender body for the step description. Use
                    "filled" when this slide is one of several lavender-only
                    slides in a row, to break monotony.
      steps: list of 2-4 dicts with: title, body
        Optional per-step fields:
          icon         : Unicode/font glyph to render INSIDE the circle instead
                         of the step number (e.g. "✓", "⚡", "→", "★").
          icon_font    : font name to render the icon with. Defaults to Calibri.
          circle_color : "navy" (default), "orange", "green" or "red".
                         For header_style="external" controls the circle fill.
                         For header_style="filled" controls the BAND color
                         (so the step header takes on the variation).
      conclusion (optional string -> wide card with orange stripe at bottom)
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    steps = slide_data.get("steps", [])[:4]
    if not steps:
        return

    header_style = slide_data.get("header_style", "external")
    conclusion = slide_data.get("conclusion", "")

    # Reserve top padding only for the external-circle style; the filled
    # header style integrates the band into the card itself.
    if header_style == "filled":
        area_top = CONTENT_TOP
    else:
        area_top = CONTENT_TOP + 0.25
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    has_conclusion = bool(conclusion)
    conclusion_h = 0.75 if has_conclusion else 0
    steps_h = area_h - conclusion_h - (CARD_GAP if has_conclusion else 0)

    n = len(steps)
    arrow_w = 0.4
    arrow_h = 0.25
    total_arrow_w = arrow_w * (n - 1)
    total_gap_w = CARD_GAP * (n - 1)
    card_w = (area_w - total_arrow_w - total_gap_w) / n

    # Numbered circle just above top of each card (external style only)
    circ_d = 0.55

    for i, step in enumerate(steps):
        cx = CONTENT_LEFT + i * (card_w + arrow_w + CARD_GAP)
        icon_text = step.get("icon", "").strip()
        icon_font = step.get("icon_font", FONT_FAMILY)
        circle_fill = resolve_circle_color(step)

        if header_style == "filled":
            # Card body covers full height, filled-header drawn on top.
            # The "circle color" choice drives the BAND color so callers
            # can mix navy/orange bands across steps.
            band_color = circle_fill
            # If band is navy, header text is white. If band is orange, text
            # stays white for contrast.
            draw_card_with_filled_header(
                slide, cx, area_top, card_w, steps_h,
                header_color=band_color,
                header_text=step.get("title", ""),
                header_text_color=WHITE,
                icon=icon_text if icon_text else str(i + 1),
                icon_font=icon_font,
                icon_circle_color=WHITE,
                body_fill=CARD_FILL,
            )
            # Body text below the header.
            ix = cx + CARD_PADDING
            iy = area_top + FILLED_HEADER_H + 0.10
            iw = card_w - 2 * CARD_PADDING
            ih = steps_h - FILLED_HEADER_H - 0.10 - CARD_PADDING
            draw_text_box(slide, ix, iy, iw, ih,
                          step.get("body", ""),
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.LEFT)
        else:
            # Original "external" style: circle floats above the card.
            step_fill = _resolve_card_fill(step)
            card = draw_card(slide, cx, area_top + circ_d / 2 + 0.05,
                             card_w, steps_h - circ_d / 2 - 0.05,
                             fill=step_fill)
            ix, iy, iw, ih = card["inner"]

            circ_cx = cx + card_w / 2
            circ_cy = area_top + circ_d / 2
            if icon_text:
                draw_circle(slide, circ_cx, circ_cy, circ_d,
                            fill=circle_fill,
                            text=icon_text, text_color=WHITE,
                            font_size=26, font_name=icon_font)
            else:
                draw_circle(slide, circ_cx, circ_cy, circ_d,
                            fill=circle_fill,
                            text=str(i + 1), text_color=WHITE, font_size=22)

            # Title (centered, bold, navy)
            draw_text_box(slide, ix, iy + 0.25, iw, 0.5,
                          step.get("title", ""),
                          font_size=FONT_CARD_TITLE, bold=True,
                          color=NAVY, align=PP_ALIGN.CENTER)
            # Body (centered, regular)
            draw_text_box(slide, ix, iy + 0.85, iw, ih - 0.85,
                          step.get("body", ""),
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.CENTER)

        # Arrow centered in the gap between this card and the next.
        # Layout is [card][gap+arrow][card][gap+arrow][card]... where the
        # combined slot between two adjacent cards is CARD_GAP + arrow_w wide.
        # Center of that slot = end-of-current-card + (CARD_GAP + arrow_w)/2;
        # arrow x = center - arrow_w/2.
        if i < n - 1:
            slot_w = CARD_GAP + arrow_w
            center_x = cx + card_w + slot_w / 2
            arrow_x = center_x - arrow_w / 2
            # Vertically center on the cards' visible region. For external
            # style, that means below the floating circle row; for filled,
            # it means inside the body region (below the header band).
            if header_style == "filled":
                ay = area_top + steps_h / 2 - arrow_h / 2
            else:
                ay = area_top + circ_d / 2 + (steps_h - circ_d / 2) / 2 - arrow_h / 2
            draw_arrow_right(slide, arrow_x, ay, arrow_w, arrow_h,
                             color=NAVY)

    # Conclusion card at bottom with orange stripe (result-style signal:
    # orange stripe + orange title, but same stripe width as other cards)
    if has_conclusion:
        c_y = area_top + steps_h + CARD_GAP
        card = draw_card(slide, CONTENT_LEFT, c_y, area_w, conclusion_h,
                         stripe_color=ORANGE)
        ix, iy, iw, ih = card["inner"]
        draw_text_box(slide, ix, iy, iw, 0.25,
                      "CONCLUSION", font_size=FONT_CARD_LABEL,
                      bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
        draw_text_box(slide, ix, iy + 0.30, iw, ih - 0.30,
                      conclusion, font_size=FONT_CARD_BODY,
                      color=NAVY, align=PP_ALIGN.LEFT)


# ---- cards_2_large ---------------------------------------------------------

def draw_cards_2_large(slide, slide_data):
    """
    2 large cards side-by-side. Each has a label header + body.
    Optionally one or both can host an image placeholder.

    slide_data keys:
      title, subtitle, use_case
      cards: list of 2 dicts with:
        label, body
        is_image (bool) - if True, draw image placeholder inside instead of body
        image_caption (optional)
        caption_below (optional italic line below the card body)
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    cards = (slide_data.get("cards", []) + [{}] * 2)[:2]

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    card_w = (area_w - CARD_GAP) / 2

    for i, spec in enumerate(cards):
        cx = CONTENT_LEFT + i * (card_w + CARD_GAP)
        stripe_color, stripe_w = _resolve_stripe_for_card(spec)
        card_fill = _resolve_card_fill(spec)
        card = draw_card(slide, cx, area_top, card_w, area_h,
                         stripe_color=stripe_color,
                         fill=card_fill,
                         stripe_width_in=stripe_w)
        ix, iy, iw, ih = card["inner"]

        # Label header — cards_2_large gets +4pt over the global title size
        # (the cards are the largest in the catalog, labels scale accordingly).
        label_color = _card_title_color(spec)
        if spec.get("label"):
            draw_text_box(slide, ix, iy, iw, 0.55,
                          spec["label"].upper(),
                          font_size=FONT_CARD_TITLE + 4, bold=True,
                          color=label_color, align=PP_ALIGN.LEFT)
        content_y = iy + 0.70
        content_h = ih - 0.70
        # Reserve space for caption below if present
        caption = spec.get("caption_below", "")
        if caption:
            content_h -= 0.30

        if spec.get("is_image"):
            draw_image_placeholder(slide, ix, content_y, iw, content_h,
                                   caption=spec.get("image_caption",
                                                    "Screenshot placeholder"))
        else:
            body = spec.get("body", "")
            # cards_2_large cards are physically the largest in the catalog
            # (~5.5" wide × ~5" tall), so they get a +4pt bump over the
            # global FONT_CARD_BODY default to fill the space and remain
            # readable at presentation distance.
            draw_text_box(slide, ix, content_y, iw, content_h, body,
                          font_size=FONT_CARD_BODY + 4, color=NAVY,
                          align=PP_ALIGN.LEFT)

        if caption:
            draw_text_box(slide, ix, iy + ih - 0.25, iw, 0.25,
                          caption, font_size=9, italic=True,
                          color=NAVY, align=PP_ALIGN.RIGHT)


# ---- card_1_hero_with_subcards --------------------------------------------

def draw_card_1_hero_with_subcards(slide, slide_data):
    """
    1 wide hero card on top (with optional pill badge) + 2x2 grid of 4 sub-cards.

    slide_data keys:
      title, subtitle, use_case
      header_style: "external" (default) draws each subcard as lavender with a
                    small bullet-circle next to the title. "filled" gives each
                    subcard a colored top band with the icon-circle LEFT of
                    the subcard title (white text on dark band), body lavender
                    below. Use "filled" when this slide is one of several
                    lavender-heavy slides in a row, to break monotony.
      hero: dict with: label, body, pill (optional pill text)
      subcards: list of 4 dicts with: title, body
        Optional per-subcard fields:
          icon         : Unicode/font glyph to render INSIDE the bullet circle
                         instead of the subcard number (e.g. "✓", "⚡", "→").
          icon_font    : font name for the icon (default Calibri).
          circle_color : "navy" (default), "orange", "green" or "red".
                         For header_style="filled" this controls the BAND color.
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    hero = slide_data.get("hero", {})
    subcards = (slide_data.get("subcards", []) + [{}] * 4)[:4]
    header_style = slide_data.get("header_style", "external")

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    hero_h = 1.20
    sub_h = area_h - hero_h - CARD_GAP

    # Hero card with orange left stripe (highlight style)
    hero_card = draw_card(slide, CONTENT_LEFT, area_top, area_w, hero_h,
                          stripe_color=ORANGE)
    ix, iy, iw, ih = hero_card["inner"]

    # Pill badge top-right
    pill_text = hero.get("pill", "")
    if pill_text:
        pill_w = 1.6
        pill_h = 0.30
        draw_pill_badge(slide,
                        ix + iw - pill_w, iy,
                        pill_w, pill_h, pill_text,
                        fill=NAVY, text_color=WHITE)
        label_max_w = iw - pill_w - 0.2
    else:
        label_max_w = iw

    # Hero label
    if hero.get("label"):
        draw_text_box(slide, ix, iy, label_max_w, 0.30,
                      hero["label"].upper(),
                      font_size=FONT_CARD_LABEL + 1, bold=True,
                      color=ORANGE, align=PP_ALIGN.LEFT)
    # Hero body
    if hero.get("body"):
        draw_text_box(slide, ix, iy + 0.35, iw, ih - 0.35,
                      hero["body"],
                      font_size=FONT_CARD_BODY, color=NAVY,
                      align=PP_ALIGN.LEFT)

    # 2x2 subcards
    sub_w = (area_w - CARD_GAP) / 2
    sub_card_h = (sub_h - CARD_GAP) / 2
    sub_top = area_top + hero_h + CARD_GAP

    positions = [
        (CONTENT_LEFT,             sub_top),
        (CONTENT_LEFT + sub_w + CARD_GAP, sub_top),
        (CONTENT_LEFT,             sub_top + sub_card_h + CARD_GAP),
        (CONTENT_LEFT + sub_w + CARD_GAP, sub_top + sub_card_h + CARD_GAP),
    ]

    for i, (spec, (sx, sy)) in enumerate(zip(subcards, positions)):
        icon_text = spec.get("icon", "").strip()
        icon_font = spec.get("icon_font", FONT_FAMILY)
        circle_fill = resolve_circle_color(spec)

        if header_style == "filled":
            # Subcard with a colored top band. Band color follows circle_color
            # so callers get the navy/orange/green/red variation visibly.
            draw_card_with_filled_header(
                slide, sx, sy, sub_w, sub_card_h,
                header_color=circle_fill,
                header_text=spec.get("title", ""),
                header_text_color=WHITE,
                icon=icon_text if icon_text else str(i + 1),
                icon_font=icon_font,
                icon_circle_color=WHITE,
                body_fill=CARD_FILL,
            )
            # Body text below the header band.
            ix2 = sx + CARD_PADDING
            iy2 = sy + FILLED_HEADER_H + 0.10
            iw2 = sub_w - 2 * CARD_PADDING
            ih2 = sub_card_h - FILLED_HEADER_H - 0.10 - CARD_PADDING
            draw_text_box(slide, ix2, iy2, iw2, ih2,
                          spec.get("body", ""),
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.LEFT)
        else:
            # Original "external" style: lavender body, small bullet circle
            # next to title.
            stripe_color, stripe_w = _resolve_stripe_for_card(spec)
            sub_fill = _resolve_card_fill(spec)
            sub = draw_card(slide, sx, sy, sub_w, sub_card_h,
                            stripe_color=stripe_color,
                            fill=sub_fill,
                            stripe_width_in=stripe_w)
            ix2, iy2, iw2, ih2 = sub["inner"]

            bd = 0.40
            if icon_text:
                draw_circle(slide, ix2 + bd / 2, iy2 + bd / 2 + 0.02,
                            bd, fill=circle_fill, text=icon_text,
                            text_color=WHITE, font_size=16,
                            font_name=icon_font)
            else:
                draw_circle(slide, ix2 + bd / 2, iy2 + bd / 2 + 0.02,
                            bd, fill=circle_fill, text=str(i + 1),
                            text_color=WHITE, font_size=16)

            title_x = ix2 + bd + 0.15
            title_w = iw2 - bd - 0.15
            draw_text_box(slide, title_x, iy2, title_w, 0.30,
                          spec.get("title", ""),
                          font_size=FONT_CARD_TITLE, bold=True,
                          color=NAVY, align=PP_ALIGN.LEFT)
            draw_text_box(slide, ix2, iy2 + 0.40, iw2, ih2 - 0.40,
                          spec.get("body", ""),
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.LEFT)


# ---- card_hero_with_white_subcards ----------------------------------------

def draw_card_hero_with_white_subcards(slide, slide_data):
    """
    1 large lavender outer container with a centered hero title + short hero
    description at the top, and 4 nested WHITE sub-cards in a 2x2 grid below.

    Replicates the "Code @ Bmatix" pattern from slide 6 of the Claude Cards
    template: the white-on-lavender inner cards create visual depth that a
    flat lavender grid cannot. Use this layout specifically as the "break the
    lavender pattern" relief after a slide that already used 4 lavender cards.

    slide_data keys:
      title, subtitle, use_case  (passed to chrome)
      hero: dict with:
        label  (small uppercase label above hero title, optional)
        title  (large hero title, centered, navy bold)
        body   (1-2 sentence hero description, centered, navy regular)
      subcards: list of 4 dicts, each with:
        icon         (modern emoji shown top-left of the white card)
        title        (bold navy card title)
        body         (1-2 line description)
        Optional:
          icon_font  (font for the icon; defaults to FONT_FAMILY)
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    hero = slide_data.get("hero", {})
    subcards = (slide_data.get("subcards", []) + [{}] * 4)[:4]

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    # Draw the large lavender outer container card filling the content area.
    # No stripe — the hero is the visual anchor, not the stripe.
    outer = draw_card(slide, CONTENT_LEFT, area_top, area_w, area_h,
                      fill=CARD_FILL)

    # Hero region: top ~25% of the outer card. Centered text.
    hero_h = 1.30
    hero_pad_x = 0.40  # extra horizontal padding inside the outer card
    hero_x = CONTENT_LEFT + hero_pad_x
    hero_w = area_w - 2 * hero_pad_x
    hero_y = area_top + 0.25

    # Optional small uppercase label above the hero title (orange, centered).
    label = hero.get("label", "")
    if label:
        draw_text_box(slide, hero_x, hero_y, hero_w, 0.25,
                      label.upper(),
                      font_size=FONT_CARD_LABEL, bold=True,
                      color=ORANGE, align=PP_ALIGN.CENTER)
        hero_title_y = hero_y + 0.30
    else:
        hero_title_y = hero_y

    # Hero title (centered, navy bold, larger than card titles).
    hero_title = hero.get("title", "")
    if hero_title:
        draw_text_box(slide, hero_x, hero_title_y, hero_w, 0.40,
                      hero_title,
                      font_size=FONT_CARD_TITLE + 4, bold=True,
                      color=NAVY, align=PP_ALIGN.CENTER)
        hero_body_y = hero_title_y + 0.45
    else:
        hero_body_y = hero_title_y

    # Hero body (centered, navy regular).
    hero_body = hero.get("body", "")
    if hero_body:
        body_h = max(0.30, (area_top + hero_h) - hero_body_y)
        draw_text_box(slide, hero_x, hero_body_y, hero_w, body_h,
                      hero_body,
                      font_size=FONT_CARD_BODY + 1, color=NAVY,
                      align=PP_ALIGN.CENTER)

    # 2x2 grid of 4 WHITE sub-cards in the remaining region.
    sub_top = area_top + hero_h + 0.20
    sub_bottom = area_bottom - 0.30  # leave breathing room inside outer card
    sub_area_h = sub_bottom - sub_top

    inset = 0.40  # how far the white grid sits inside the lavender outer card
    sub_grid_left = CONTENT_LEFT + inset
    sub_grid_w = area_w - 2 * inset
    sub_gap = CARD_GAP

    sub_w = (sub_grid_w - sub_gap) / 2
    sub_h = (sub_area_h - sub_gap) / 2

    positions = [
        (sub_grid_left,                 sub_top),
        (sub_grid_left + sub_w + sub_gap, sub_top),
        (sub_grid_left,                 sub_top + sub_h + sub_gap),
        (sub_grid_left + sub_w + sub_gap, sub_top + sub_h + sub_gap),
    ]

    for spec, (sx, sy) in zip(subcards, positions):
        # White inner card — no stripe, white fill.
        inner = draw_card(slide, sx, sy, sub_w, sub_h, fill=WHITE)
        ix, iy, iw, ih = inner["inner"]

        icon_text = spec.get("icon", "").strip()
        icon_font = spec.get("icon_font", FONT_FAMILY)

        # Emoji icon top-left (no circle — just the glyph, larger size).
        if icon_text:
            draw_text_box(slide, ix, iy, 0.50, 0.40,
                          icon_text,
                          font_size=22, color=NAVY,
                          font_name=icon_font, align=PP_ALIGN.LEFT)
            title_x = ix + 0.55
            title_w = iw - 0.55
        else:
            title_x = ix
            title_w = iw

        # Card title (bold navy, right of the icon).
        draw_text_box(slide, title_x, iy + 0.05, title_w, 0.35,
                      spec.get("title", ""),
                      font_size=FONT_CARD_TITLE, bold=True,
                      color=NAVY, align=PP_ALIGN.LEFT)

        # Card body (below the icon+title row, full width).
        body_y = iy + 0.50
        body_h = ih - 0.50
        draw_text_box(slide, ix, body_y, iw, body_h,
                      spec.get("body", ""),
                      font_size=FONT_CARD_BODY, color=NAVY,
                      align=PP_ALIGN.LEFT)


# ---- cards_2_bulletlist_plus_wide -----------------------------------------

def draw_cards_2_bulletlist_plus_wide(slide, slide_data):
    """
    2 bullet-list cards side-by-side on top + 1 wide example card below.

    slide_data keys:
      title, subtitle, use_case
      cards: list of 2 dicts with: header, header_color ('navy'|'orange'),
             bullets (list of strings)
      example (optional dict): label, body
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""),
                      use_case=slide_data.get("use_case", ""))

    cards = (slide_data.get("cards", []) + [{}] * 2)[:2]
    example = slide_data.get("example")

    area_top = CONTENT_TOP
    area_bottom = CONTENT_BOTTOM - 0.3
    area_h = area_bottom - area_top
    area_w = CONTENT_RIGHT - CONTENT_LEFT

    has_example = bool(example)
    example_h = 1.20 if has_example else 0
    top_h = area_h - example_h - (CARD_GAP if has_example else 0)

    card_w = (area_w - CARD_GAP) / 2

    for i, spec in enumerate(cards):
        cx = CONTENT_LEFT + i * (card_w + CARD_GAP)
        card_fill = _resolve_card_fill(spec)
        stripe_color, stripe_w = _resolve_stripe_for_card(spec)
        card = draw_card(slide, cx, area_top, card_w, top_h,
                         stripe_color=stripe_color,
                         fill=card_fill,
                         stripe_width_in=stripe_w)
        ix, iy, iw, ih = card["inner"]

        header_color = (ORANGE if spec.get("header_color", "").lower() == "orange"
                        else NAVY)
        if spec.get("header"):
            draw_text_box(slide, ix, iy, iw, 0.30,
                          spec["header"],
                          font_size=FONT_CARD_TITLE, bold=True,
                          color=header_color, align=PP_ALIGN.LEFT)

        # Bullets joined by newlines, prefixed with •
        bullets = spec.get("bullets", [])
        body = "\n".join(f"•  {b}" for b in bullets)
        if body:
            draw_text_box(slide, ix, iy + 0.40, iw, ih - 0.40,
                          body, font_size=FONT_CARD_BODY,
                          color=NAVY, align=PP_ALIGN.LEFT)

    if has_example:
        ex_y = area_top + top_h + CARD_GAP
        # The example block at the bottom uses the beige/sand fill since it's
        # showcasing a concrete example/demo, not real content. Stripe uses the
        # tan/brown EXAMPLE color (slightly darker than the beige fill) to
        # match the auto-styled example cards elsewhere in the deck.
        card = draw_card(slide, CONTENT_LEFT, ex_y, area_w, example_h,
                         stripe_color=STRIPE_EXAMPLE,
                         fill=CARD_FILL_EXAMPLE)
        ix, iy, iw, ih = card["inner"]
        if example.get("label"):
            draw_text_box(slide, ix, iy, iw, 0.30,
                          example["label"].upper(),
                          font_size=FONT_CARD_LABEL, bold=True,
                          color=NAVY, align=PP_ALIGN.LEFT)
        if example.get("body"):
            draw_text_box(slide, ix, iy + 0.35, iw, ih - 0.35,
                          example["body"],
                          font_size=FONT_CARD_BODY, color=NAVY,
                          align=PP_ALIGN.LEFT)


def draw_decision_tree(slide, slide_data):
    """
    Decision-tree layout: a central root oval branches into 2-4 paths,
    each path containing a sequence of 2-5 numbered step cards.

    Visual structure (matches slide 8 of Claude_Cards_template.pptx):
      - Central navy oval at the top, containing an emoji (the "root")
      - One short navy decision-question line below the oval
      - Horizontal navy connector line spanning the path columns
      - For each path:
          * Navy banner header at top with an emoji and a label (white text)
          * 2-5 white step-cards stacked vertically below; each card has a
            small navy oval with a white step-number on the left and a short
            label on the right (2-6 words)
          * Thin navy vertical connectors between consecutive step cards

    slide_data keys:
      title, subtitle              (rendered via native placeholders)
      root: {icon: "🤖", label: "Question text?"}    (emoji + decision question)
      branches: list of 2-4 dicts, each with:
        icon       : emoji shown white inside the banner header
        label      : banner-header text (white, e.g. "Fast path · No chat")
        steps      : list of 2-5 short strings, OR list of dicts with key
                     "label" (and optionally "icon"). Step circles always
                     stay navy with white numerals — no per-step color mix.

    Design rules (locked):
      - Banner headers and root oval are ALWAYS navy. No per-branch color mix.
      - Step ovals are ALWAYS navy with white numerals (1, 2, 3, ...).
      - 5+ branches is rejected (would be too cramped) — caller should split.
      - Step text is kept short (2-6 words); long text is truncated visually
        by the card width but no auto-wrap shrinking happens here.
    """
    draw_slide_chrome(slide,
                      title=slide_data["title"],
                      subtitle=slide_data.get("subtitle", ""))

    root = slide_data.get("root", {}) or {}
    branches = slide_data.get("branches", []) or []
    if not branches:
        return
    if len(branches) > 4:
        branches = branches[:4]

    n = len(branches)

    # --- Layout geometry --------------------------------------------------
    # Available content zone: CONTENT_LEFT..CONTENT_RIGHT, CONTENT_TOP..CONTENT_BOTTOM
    area_w = CONTENT_RIGHT - CONTENT_LEFT
    area_top = CONTENT_TOP

    # Root oval centered horizontally at the top of the content area.
    root_d = 0.72
    root_cx = CONTENT_LEFT + area_w / 2
    root_cy = area_top + root_d / 2 - 0.05

    # Question text just under the root oval
    question_y = root_cy + root_d / 2 + 0.05
    question_h = 0.30

    # Horizontal connector line (just under the question text)
    h_line_y = question_y + question_h + 0.08

    # Banner header band
    banner_top = h_line_y + 0.25
    banner_h = 0.45

    # Vertical short stubs from h_line down to each banner (drawn as thin rects)
    stub_top = h_line_y
    stub_bottom = banner_top

    # Step-card geometry
    step_card_h = 0.55
    step_card_gap = 0.18
    step_area_top = banner_top + banner_h + 0.18

    # Per-branch column width
    # Use 0.30" inter-column gap so 4 branches don't crowd
    col_gap = 0.30
    total_gap = col_gap * (n - 1)
    col_w = (area_w - total_gap) / n

    def col_x(i):
        return CONTENT_LEFT + i * (col_w + col_gap)

    # --- Draw root oval --------------------------------------------------
    root_icon = (root.get("icon") or "").strip()
    draw_circle(slide, root_cx, root_cy, root_d,
                fill=NAVY, text=root_icon,
                font_size=24, font_name=FONT_FAMILY)

    # --- Draw question text under root oval ------------------------------
    question_label = (root.get("label") or "").strip()
    if question_label:
        q_w = min(6.0, area_w)
        q_x = root_cx - q_w / 2
        draw_text_box(slide, q_x, question_y, q_w, question_h,
                      question_label,
                      font_size=12, bold=True, color=NAVY,
                      align=PP_ALIGN.CENTER)

    # --- Horizontal connector line ---------------------------------------
    # Only spans from the center of the first column to the center of the last
    if n >= 2:
        first_cx = col_x(0) + col_w / 2
        last_cx = col_x(n - 1) + col_w / 2
        line_w = last_cx - first_cx
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(first_cx), Inches(h_line_y),
            Inches(line_w), Inches(0.025)
        )
        _set_solid_fill(line_shape, NAVY)
        _set_no_line(line_shape)
    # For single-branch decision tree (rare), no horizontal line — straight vertical.

    # --- Per-branch drawing -----------------------------------------------
    for i, branch in enumerate(branches):
        cx0 = col_x(i)
        col_center_x = cx0 + col_w / 2

        # Vertical stub from h_line down to banner top (thin navy rect)
        stub_h = stub_bottom - stub_top
        if stub_h > 0:
            stub_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_center_x - 0.0125), Inches(stub_top),
                Inches(0.025), Inches(stub_h)
            )
            _set_solid_fill(stub_shape, NAVY)
            _set_no_line(stub_shape)

        # Small downward-pointing arrow tip just above the banner
        tip_h = 0.10
        tip_w = 0.14
        tip = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(col_center_x - tip_w / 2),
            Inches(stub_bottom - tip_h - 0.02),
            Inches(tip_w), Inches(tip_h)
        )
        _set_solid_fill(tip, NAVY)
        _set_no_line(tip)

        # Banner header (navy, white text + emoji)
        banner = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx0), Inches(banner_top),
            Inches(col_w), Inches(banner_h)
        )
        _set_solid_fill(banner, NAVY)
        _set_no_line(banner)
        try:
            banner.adjustments[0] = 0.20
        except Exception:
            pass

        # Banner content: emoji at left, label centered
        banner_icon = (branch.get("icon") or "").strip()
        banner_label = (branch.get("label") or "").strip()
        # Render emoji + label in one text frame, emoji bigger
        if banner_icon and banner_label:
            banner_text = f"{banner_icon}  {banner_label}"
        else:
            banner_text = banner_icon or banner_label
        if banner_text:
            tf = banner.text_frame
            tf.margin_left = Inches(0.10)
            tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            # Clear any existing text
            for r in list(p.runs):
                r.text = ""
            run = p.add_run()
            run.text = banner_text
            run.font.name = FONT_FAMILY
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = WHITE

        # --- Draw step cards under banner ---------------------------------
        steps = branch.get("steps", []) or []
        steps = steps[:5]
        for s_idx, step in enumerate(steps):
            # Allow step to be either a string or a dict {label, icon}
            if isinstance(step, dict):
                step_label = (step.get("label") or "").strip()
                step_icon = (step.get("icon") or "").strip()
            else:
                step_label = str(step).strip()
                step_icon = ""

            step_y = step_area_top + s_idx * (step_card_h + step_card_gap)

            # Vertical connector between consecutive step cards
            if s_idx > 0:
                conn_top = step_y - step_card_gap + 0.04
                conn_h = step_card_gap - 0.08
                if conn_h > 0:
                    conn = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(col_center_x - 0.0125), Inches(conn_top),
                        Inches(0.025), Inches(conn_h)
                    )
                    _set_solid_fill(conn, NAVY)
                    _set_no_line(conn)

            # White step card with thin gray border
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx0), Inches(step_y),
                Inches(col_w), Inches(step_card_h)
            )
            _set_solid_fill(card, WHITE)
            _set_line(card, RGBColor(0xC7, 0xCA, 0xD8), width_pt=1.0)
            try:
                card.adjustments[0] = 0.22
            except Exception:
                pass

            # Navy circle with step number on the left
            circ_d = 0.40
            circ_cx = cx0 + 0.05 + circ_d / 2
            circ_cy = step_y + step_card_h / 2
            number_text = step_icon if step_icon else str(s_idx + 1)
            draw_circle(slide, circ_cx, circ_cy, circ_d,
                        fill=NAVY, text=number_text,
                        font_size=12, font_name=FONT_FAMILY)

            # Step label text to the right of the circle
            if step_label:
                txt_x = cx0 + 0.05 + circ_d + 0.10
                txt_w = col_w - (0.05 + circ_d + 0.10) - 0.10
                draw_text_box(slide, txt_x, step_y + 0.05,
                              txt_w, step_card_h - 0.10,
                              step_label,
                              font_size=11, bold=True,
                              color=RGBColor(0x1A, 0x1A, 0x1A),
                              align=PP_ALIGN.LEFT)


# =============================================================================
# CARD LAYOUT DISPATCH
# =============================================================================

CARD_LAYOUTS = {
    "cards_2x2_compare":             draw_cards_2x2_compare,
    "cards_left_image_right":        draw_cards_left_image_right,
    "cards_numbered_steps":          draw_cards_numbered_steps,
    "cards_2_large":                 draw_cards_2_large,
    "card_1_hero_with_subcards":     draw_card_1_hero_with_subcards,
    "card_hero_with_white_subcards": draw_card_hero_with_white_subcards,
    "cards_2_bulletlist_plus_wide":  draw_cards_2_bulletlist_plus_wide,
    "decision_tree":                 draw_decision_tree,
}


def draw_card_slide(slide, slide_data):
    """Dispatch to the correct card-layout drawer based on slide_data['layout']."""
    layout = slide_data.get("layout")
    fn = CARD_LAYOUTS.get(layout)
    if fn is None:
        raise ValueError(
            f"Unknown card layout '{layout}'. "
            f"Available: {sorted(CARD_LAYOUTS.keys())}"
        )
    fn(slide, slide_data)


# =============================================================================
# INTRO + CLOSING FILLERS (template slides 27 and 51)
# =============================================================================

def _fill_placeholder_by_type(slide, ph_type_name: str, new_text: str,
                              min_height_in: float = 0.20):
    """Find the first placeholder of the given type on a slide whose height
    is at least min_height_in inches, and replace its text. Returns True if
    found and replaced, False otherwise.

    The `min_height_in` filter exists because the Bmatix templates contain
    "phantom" BODY placeholders (tiny ~0.14" tall, empty, geometrically out
    of the way) that PowerPoint never renders visibly. Without filtering by
    height, _fill_placeholder_by_type would write into the first BODY it
    sees — which is usually the phantom, not the visible subtitle/body.
    Setting min_height_in to 0.20" skips the phantoms (typically 0.14") and
    selects only visible placeholders (subtitle ≥ 0.46", body ≥ 0.30").

    Type matching is by string (e.g. 'CENTER_TITLE (3)', 'SUBTITLE (4)',
    'TITLE (1)', 'BODY (2)', 'FOOTER (15)').
    """
    for shape in slide.shapes:
        if not shape.is_placeholder or not shape.has_text_frame:
            continue
        try:
            current_type = str(shape.placeholder_format.type)
        except Exception:
            continue
        if current_type != ph_type_name:
            continue
        # Filter out phantom placeholders (tiny height)
        h_in = (shape.height or 0) / 914400
        if h_in < min_height_in:
            continue
        _set_shape_text_inplace(shape, new_text)
        return True
    return False


def fill_intro_slide(slide, *, title: str, subtitle: str = "",
                     author_date: str = ""):
    """Fill the intro slide (slide 1 of either template).

    Both templates use the "Title light" layout, which has:
      - CENTER_TITLE placeholder (the big "Title")
      - SUBTITLE placeholder (the smaller orange/colored subline)

    `author_date` is unused on the intro slide of these templates — there
    is no name/date placeholder on the intro. It's accepted for API
    symmetry with fill_closing_slide.
    """
    if title:
        _fill_placeholder_by_type(slide, "CENTER_TITLE (3)", title)
    if subtitle:
        _fill_placeholder_by_type(slide, "SUBTITLE (4)", subtitle)
    # author_date intentionally ignored — intro slide has no such placeholder
    _ = author_date


def fill_closing_slide(slide, *, author_date: str = ""):
    """Fill the closing slide (last slide of either template).

    Both templates use the "End Slide" layout, which has:
      - CENTER_TITLE = "Questions?" (left as-is)
      - SUBTITLE placeholder = 'Presented by "Presenter" on "Date"'
        (we replace with the presenter + date string)

    The title "Questions?" stays. Only the subtitle gets personalized.
    """
    if author_date:
        # Format: "Presented by <Name> on <Date>"
        # The template uses curly quotes around Presenter/Date as placeholders,
        # but we substitute the entire string for clarity.
        subtitle_text = f"Presented by {author_date}" if author_date.lower().startswith("presented") is False else author_date
        # Actually, simpler — accept author_date as the full string the caller
        # wants to appear; let the caller decide format.
        _fill_placeholder_by_type(slide, "SUBTITLE (4)", author_date)


# =============================================================================
# LOW-LEVEL TEXT-REPLACEMENT PRIMITIVE
# =============================================================================

def _set_shape_text_inplace(shape, new_text: str):
    """
    Replace a shape's text while preserving the formatting of the first run.

    This is the primitive used by _fill_placeholder_by_type and by the
    intro/closing fillers. It overwrites text content but leaves all the
    paragraph-level and run-level formatting attributes (font, size, color,
    bold, alignment) from the FIRST run of the FIRST paragraph in place.
    All subsequent runs and paragraphs are stripped — useful for replacing
    multi-line placeholder text ("Level\nLevel\nLevel...") with a single
    clean string.

    If the placeholder has no runs (empty), falls back to assigning to
    tf.text directly, which uses default formatting from the layout/master.
    """
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        p_elem = tf.paragraphs[0]._p
        runs_xml = p_elem.findall(qn("a:r"))
        for r in runs_xml[1:]:
            p_elem.remove(r)
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
    else:
        tf.text = new_text
